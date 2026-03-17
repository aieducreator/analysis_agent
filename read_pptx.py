import sys
from pptx import Presentation

def extract_text(pptx_path):
    prs = Presentation(pptx_path)
    with open('pptx_content.txt', 'w', encoding='utf-8') as f:
        for slide_idx, slide in enumerate(prs.slides):
            f.write(f"--- Slide {slide_idx + 1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    f.write(shape.text + "\n")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        extract_text(sys.argv[1])
    else:
        print("Provide pptx path")
